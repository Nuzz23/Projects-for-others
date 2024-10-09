from pptx import *
from pptx.util import Pt

# Ricreare la presentazione PowerPoint con la struttura fornita
prs = Presentation()

# Funzione per aggiungere una slide con titolo e contenuto
def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Titolo e contenuto
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = content
    for paragraph in content_placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)

# Slide 1: Titolo
slide_layout = prs.slide_layouts[0]  # Solo titolo
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Laboratorio di Acustica"
subtitle.text = ("Valutazione isolamento acustico in edifici e di elementi di edificio\n"
                 "Corso: Fisica dell’edificio e climatizzazione\n"
                 "Anno accademico: 2023/2024\n"
                 "Data: 2/04/24\n"
                 "Autori: Matteo Lopresti S294512, Andrea Leggio S299545")

# Slide 2: Indice
add_slide(prs, "Indice",
          "1. Introduzione\n"
          "2. Raccolta Dati\n"
          "3. Misurazione per l’isolamento per via aerea [Parte 1]\n"
          "4. Misurazione per l’isolamento per calpestio [Parte 2]\n"
          "5. Calcoli\n"
          "6. Conclusioni")

# Slide 3: Introduzione
add_slide(prs, "Introduzione",
          "Relazione sviluppata alla conclusione di un percorso formativo sull’acustica.\n"
          "Analisi delle caratteristiche sonore di un edificio:\n"
          "- Isolamento acustico per via aerea (Parte 1)\n"
          "- Isolamento del rumore per calpestio (Parte 2)\n"
          "Obiettivo: valutare se gli ambienti sono a norma.")

# Slide 4: Raccolta Dati
add_slide(prs, "Raccolta Dati",
          "Strumenti utilizzati:\n"
          "- Sorgente sonora (cassa bluetooth, stereo, impianti audio hi-fi)\n"
          "- Smartphone Android con app 'Sound Level Meter' e 'USB Reverberation Meter'\n"
          "- Microsoft Excel\n"
          "- 4 spaghetti crudi per standardizzare le misurazioni\n"
          "- Libro di massa 2-3 kg per misurare il tempo di riverberazione\n"
          "Edifici:\n"
          "- Edificio A: isolamento per via aerea\n"
          "- Edificio B: isolamento per calpestio")

# Slide 5: Misurazione per l’isolamento per via aerea [Parte 1]
add_slide(prs, "Misurazione per l’isolamento per via aerea [Parte 1]",
          "Edificio A a Vignone (VB), condominio a due piani\n"
          "Data delle misurazioni: 29/03/2024\n"
          "Strumenti: cassa bluetooth LGPK3W, dizionario italiano-inglese Garzanti\n"
          "Ambiente ricevente: cucina (23,87 m³, superficie parete 9,55 m²)\n"
          "Ambiente emittente: sala da pranzo (55,45 m³)")

# Slide 6: Risultati Misurazione per l’isolamento per via aerea
add_slide(prs, "Risultati Misurazione per l’isolamento per via aerea",
          "Taratura dello strumento:\n"
          "- Media valore di picco: 31,95 dB\n"
          "- Correzione: 45,35 dB\n"
          "Misurazioni pressione sonora:\n"
          "- Livello globale di fondo: 73,8 dB\n"
          "- Livello con rumore rosa: 82 dB\n"
          "Tabelle dei dati sperimentali di pressione sonora e tempo di riverberazione")

# Slide 7: Misurazione per l’isolamento per calpestio [Parte 2]
add_slide(prs, "Misurazione per l’isolamento per calpestio [Parte 2]",
          "Edificio B a Ragusa (RG), villetta a schiera\n"
          "Data delle misurazioni: 03/04/2024\n"
          "Strumenti: rulletta metrica, dizionario devoto-oli, spaghetti per calibrazione\n"
          "Ambiente ricevente: cucina (volume totale 56,58 m³)\n"
          "Ambiente disturbante: stanza da letto (volume totale 37,74 m³)")

# Slide 8: Risultati Misurazione per l’isolamento per calpestio
add_slide(prs, "Risultati Misurazione per l’isolamento per calpestio",
          "Taratura dello strumento:\n"
          "- Media valore di picco: 18,34 dB\n"
          "- Correzione: 58,96 dB\n"
          "Tabelle dei dati sperimentali di pressione sonora e tempo di riverberazione")

# Slide 9: Calcoli
add_slide(prs, "Calcoli",
          "Confronto con i valori del D.P.C.M del 1997\n"
          "Parte 1: indice di isolamento acustico apparente a parete di 24 (R’)\n"
          "Parte 2: indice di valutazione pari a 69\n"
          "Grafici dei risultati ottenuti")

# Slide 10: Conclusioni
add_slide(prs, "Conclusioni",
          "Parte 1: isolamento acustico apparente al di sotto dei limiti normativi\n"
          "Parte 2: livello di rumore di calpestio supera il limite normativo\n"
          "Importanza dell’esperienza e della precisione strumentale\n"
          "Risultati non scontati data l’inesperienza\n"
          "Esperienza migliorabile con misurazioni in ambienti a norma")

# Salvare la presentazione
file_path = "Laboratorio_di_acustica_presentazione_finale.pptx"
prs.save(file_path)

file_path